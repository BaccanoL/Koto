#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文档内容读取器 - 读取并结构化各类文档
支持：PPT、Word、Excel
"""

import os
from typing import Dict, List, Any, Optional
from pathlib import Path


class DocumentReader:
    """文档内容读取器"""
    
    @staticmethod
    def read_ppt(file_path: str) -> Dict[str, Any]:
        """
        读取PPT文件并提取结构化内容
        
        Returns:
            {
                "type": "ppt",
                "title": "演示标题",
                "slide_count": 10,
                "slides": [
                    {
                        "index": 0,
                        "title": "幻灯片标题",
                        "content": ["要点1", "要点2"],
                        "notes": "演讲者备注"
                    },
                    ...
                ]
            }
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_data = []
            
            for idx, slide in enumerate(prs.slides):
                slide_info = {
                    "index": idx,
                    "title": "",
                    "content": [],
                    "notes": ""
                }
                
                # 提取标题
                if slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text
                
                # 提取内容（文本框）
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        # 排除标题
                        if shape != slide.shapes.title:
                            # 尝试提取列表项
                            if hasattr(shape, "text_frame"):
                                for paragraph in shape.text_frame.paragraphs:
                                    text = paragraph.text.strip()
                                    if text:
                                        slide_info["content"].append(text)
                            else:
                                slide_info["content"].append(shape.text)
                
                # 提取备注
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        slide_info["notes"] = notes_slide.notes_text_frame.text
                
                slides_data.append(slide_info)
            
            return {
                "success": True,
                "type": "ppt",
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "slide_count": len(slides_data),
                "slides": slides_data
            }
            
        except ImportError:
            return {
                "success": False,
                "error": "需要安装 python-pptx: pip install python-pptx"
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"读取PPT失败: {str(e)}"
            }
    
    @staticmethod
    def read_word(file_path: str) -> Dict[str, Any]:
        """
        读取Word文件并提取结构化内容
        
        Returns:
            {
                "type": "word",
                "paragraphs": [
                    {"type": "heading", "level": 1, "text": "标题"},
                    {"type": "paragraph", "text": "正文内容"},
                    {"type": "list", "text": "列表项"},
                    ...
                ],
                "tables": [...]
            }
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs_data = []
            tables_data = []
            
            # 提取段落（增强：包含格式信息）
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                
                # 提取runs的格式信息
                runs_info = []
                has_format_change = False
                for run in para.runs:
                    if not run.text:
                        continue
                    run_data = {
                        "text": run.text,
                        "bold": run.bold if run.bold is not None else False,
                        "italic": run.italic if run.italic is not None else False,
                    }
                    if run.font.size:
                        run_data["size_pt"] = run.font.size.pt
                    if run.font.color and run.font.color.rgb:
                        run_data["color"] = str(run.font.color.rgb)
                        has_format_change = True
                    if run.bold or run.italic:
                        has_format_change = True
                    runs_info.append(run_data)
                
                para_info = {
                    "text": para.text,
                    "style": para.style.name if para.style else "Normal",
                    "runs": runs_info,  # 新增：run级别的格式信息
                    "has_format_change": has_format_change,  # 新增：是否有格式变化
                }
                
                # 检测标题
                if para.style and "Heading" in para.style.name:
                    para_info["type"] = "heading"
                    para_info["level"] = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                # 检测列表
                elif para.style and "List" in para.style.name:
                    para_info["type"] = "list"
                else:
                    para_info["type"] = "paragraph"
                
                paragraphs_data.append(para_info)
            
            # 提取表格
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables_data.append(table_data)
            
            return {
                "success": True,
                "type": "word",
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "paragraph_count": len(paragraphs_data),
                "paragraphs": paragraphs_data,
                "table_count": len(tables_data),
                "tables": tables_data
            }
            
        except ImportError:
            return {
                "success": False,
                "error": "需要安装 python-docx: pip install python-docx"
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"读取Word失败: {str(e)}"
            }
    
    @staticmethod
    def read_excel(file_path: str) -> Dict[str, Any]:
        """
        读取Excel文件并提取结构化内容
        
        Returns:
            {
                "type": "excel",
                "sheets": [
                    {
                        "name": "Sheet1",
                        "rows": [[cell1, cell2, ...], ...],
                        "row_count": 10,
                        "col_count": 5
                    },
                    ...
                ]
            }
        """
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            sheets_data = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    # 跳过全空行
                    if all(cell is None or str(cell).strip() == "" for cell in row):
                        continue
                    rows_data.append(list(row))
                
                if rows_data:
                    sheets_data.append({
                        "name": sheet_name,
                        "rows": rows_data,
                        "row_count": len(rows_data),
                        "col_count": len(rows_data[0]) if rows_data else 0
                    })
            
            return {
                "success": True,
                "type": "excel",
                "file_path": file_path,
                "file_name": os.path.basename(file_path),
                "sheet_count": len(sheets_data),
                "sheets": sheets_data
            }
            
        except ImportError:
            return {
                "success": False,
                "error": "需要安装 openpyxl: pip install openpyxl"
            }
        except Exception as e:
            return {
                "success": False,
                "error": f"读取Excel失败: {str(e)}"
            }
    
    @staticmethod
    def read_document(file_path: str) -> Dict[str, Any]:
        """
        自动识别文件类型并读取
        
        Args:
            file_path: 文件路径
        
        Returns:
            结构化的文档内容
        """
        ext = Path(file_path).suffix.lower()
        
        if ext in ['.ppt', '.pptx']:
            return DocumentReader.read_ppt(file_path)
        elif ext in ['.doc', '.docx']:
            return DocumentReader.read_word(file_path)
        elif ext in ['.xls', '.xlsx']:
            return DocumentReader.read_excel(file_path)
        else:
            return {
                "success": False,
                "error": f"不支持的文件格式: {ext}"
            }
    
    @staticmethod
    def format_for_ai(doc_data: Dict[str, Any]) -> str:
        """
        将文档内容格式化为适合AI分析的文本
        
        Args:
            doc_data: read_document返回的结构化数据
        
        Returns:
            格式化的Markdown文本
        """
        if not doc_data.get("success"):
            return f"错误: {doc_data.get('error', '未知错误')}"
        
        doc_type = doc_data.get("type")
        output = []
        
        output.append(f"# 文档分析\n")
        output.append(f"**文件名**: {doc_data.get('file_name', 'unknown')}")
        output.append(f"**文件类型**: {doc_type.upper()}\n")
        
        if doc_type == "ppt":
            output.append(f"**幻灯片总数**: {doc_data.get('slide_count', 0)}\n")
            output.append("## 幻灯片内容\n")
            
            for slide in doc_data.get("slides", []):
                output.append(f"### 幻灯片 {slide['index'] + 1}: {slide['title'] or '(无标题)'}")
                
                if slide['content']:
                    output.append("**内容**:")
                    for item in slide['content']:
                        output.append(f"- {item}")
                
                if slide['notes']:
                    output.append(f"**备注**: {slide['notes']}")
                
                output.append("")  # 空行

        elif doc_type == "word":
            output.append(f"**段落总数**: {doc_data.get('paragraph_count', 0)}")
            output.append(f"**表格总数**: {doc_data.get('table_count', 0)}\n")
            output.append("## 文档内容\n")
            
            for para in doc_data.get("paragraphs", []):
                if para['type'] == 'heading':
                    level = para.get('level', 1)
                    output.append(f"{'#' * (level + 2)} {para['text']}")
                elif para['type'] == 'list':
                    output.append(f"- {para['text']}")
                else:
                    # 新增：如果有格式变化，标记出来
                    if para.get('has_format_change') and para.get('runs'):
                        formatted_text = ""
                        for run in para['runs']:
                            text = run['text']
                            if run.get('bold'):
                                text = f"**{text}**"
                            if run.get('italic'):
                                text = f"*{text}*"
                            if run.get('color'):
                                text = f"[{text}](颜色:{run['color']})"
                            formatted_text += text
                        output.append(formatted_text + "  ← [此段落有格式变化，包含粗体/斜体/颜色等]")
                    else:
                        output.append(para['text'])
                output.append("")
        
        elif doc_type == "excel":
            output.append(f"**工作表总数**: {doc_data.get('sheet_count', 0)}\n")
            
            for sheet in doc_data.get("sheets", []):
                output.append(f"## 工作表: {sheet['name']}")
                output.append(f"**尺寸**: {sheet['row_count']}行 × {sheet['col_count']}列\n")
                
                # 只显示前10行
                rows_to_show = sheet['rows'][:10]
                for row in rows_to_show:
                    output.append("| " + " | ".join(str(cell) if cell else "" for cell in row) + " |")
                
                if sheet['row_count'] > 10:
                    output.append(f"\n...（还有{sheet['row_count'] - 10}行）\n")
        
        return "\n".join(output)

    @staticmethod
    def _extract_text_from_doc_data(doc_data: Dict[str, Any]) -> str:
        """将结构化文档数据提取为纯文本（用于 TXT 转换）"""
        if not doc_data or not doc_data.get("success"):
            return ""

        lines = []
        for para in doc_data.get("paragraphs", []):
            text = (para.get("text") or "").strip()
            if text:
                lines.append(text)

        # 追加表格文本（按行拼接）
        for table in doc_data.get("tables", []):
            for row in table:
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row]).strip()
                if row_text:
                    lines.append(row_text)

        return "\n".join(lines)


if __name__ == "__main__":
    # 测试
    reader = DocumentReader()
    
    # 测试读取PPT
    test_ppt = r"workspace\documents\MicroLED技术全景解析_20260201_194320.pptx"
    if os.path.exists(test_ppt):
        print("=" * 60)
        print("测试读取PPT")
        print("=" * 60)
        result = reader.read_ppt(test_ppt)
        print(f"成功: {result.get('success')}")
        print(f"幻灯片数: {result.get('slide_count')}")
        
        formatted = reader.format_for_ai(result)
        print("\n格式化输出:")
        print(formatted[:500] + "..." if len(formatted) > 500 else formatted)
