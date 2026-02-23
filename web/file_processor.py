#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
æ–‡ä»¶å¤„ç†å™¨ - ç»Ÿä¸€å¤„ç†å„ç§æ–‡ä»¶ç±»å‹çš„è§£æå’Œå†…å®¹æå–
æ”¯æŒæ ¼å¼: .docx, .pdf, .xlsx, .txt, å›¾ç‰‡ç­‰
"""

import os
import mimetypes
from typing import Dict, Any, Optional, Tuple


class FileProcessor:
    """æ–‡ä»¶å¤„ç†å™¨ - æå–æ–‡ä»¶å†…å®¹å’Œå…ƒæ•°æ®"""
    
    @staticmethod
    def process_file(filepath: str) -> Dict[str, Any]:
        """
        å¤„ç†æ–‡ä»¶ï¼Œè¿”å›å¤„ç†ç»“æœ
        
        Args:
            filepath: æ–‡ä»¶è·¯å¾„
            
        Returns:
            {
                'success': bool,
                'mime_type': str,
                'filename': str,
                'text_content': str,  # æå–çš„æ–‡æœ¬å†…å®¹
                'binary_data': bytes,  # äºŒè¿›åˆ¶æ•°æ®ï¼ˆå›¾ç‰‡/PDFï¼‰
                'error': str,  # é”™è¯¯ä¿¡æ¯
                'metadata': dict  # é¢å¤–å…ƒæ•°æ®
            }
        """
        filename = os.path.basename(filepath)
        mime_type, _ = mimetypes.guess_type(filepath)
        if not mime_type:
            mime_type = "application/octet-stream"
        
        result = {
            'success': False,
            'mime_type': mime_type,
            'filename': filename,
            'text_content': '',
            'binary_data': None,
            'error': '',
            'metadata': {}
        }
        
        try:
            # æ ¹æ®æ–‡ä»¶ç±»å‹åˆ†å‘å¤„ç†
            if mime_type.startswith('image'):
                return FileProcessor._process_image(filepath, result)
            elif filename.endswith('.pdf') or mime_type == 'application/pdf':
                return FileProcessor._process_pdf(filepath, result)
            elif filename.endswith(('.doc', '.docx')):
                return FileProcessor._process_word(filepath, result)
            elif filename.endswith(('.ppt', '.pptx')):
                return FileProcessor._process_powerpoint(filepath, result)
            elif filename.endswith(('.xls', '.xlsx')):
                return FileProcessor._process_excel(filepath, result)
            elif mime_type.startswith('text') or filename.endswith('.txt'):
                return FileProcessor._process_text(filepath, result)
            else:
                # æœªçŸ¥ç±»å‹ï¼Œå°è¯•ä½œä¸ºæ–‡æœ¬è¯»å–
                return FileProcessor._process_text(filepath, result)
                
        except Exception as e:
            result['error'] = f"å¤„ç†æ–‡ä»¶å¤±è´¥: {str(e)}"
            result['success'] = False
            return result
    
    @staticmethod
    def _process_image(filepath: str, result: Dict) -> Dict:
        """å¤„ç†å›¾ç‰‡æ–‡ä»¶"""
        try:
            with open(filepath, 'rb') as f:
                result['binary_data'] = f.read()
            
            # è·å–å›¾ç‰‡å°ºå¯¸ï¼ˆå¯é€‰ï¼‰
            try:
                from PIL import Image
                with Image.open(filepath) as img:
                    result['metadata']['dimensions'] = f"{img.width}x{img.height}"
                    result['metadata']['format'] = img.format
            except ImportError:
                pass  # PILæœªå®‰è£…ï¼Œè·³è¿‡å…ƒæ•°æ®
            
            result['success'] = True
            print(f"[FileProcessor] æˆåŠŸå¤„ç†å›¾ç‰‡: {result['filename']}")
            return result
            
        except Exception as e:
            result['error'] = f"è¯»å–å›¾ç‰‡å¤±è´¥: {str(e)}"
            return result
    
    @staticmethod
    def _process_pdf(filepath: str, result: Dict) -> Dict:
        """å¤„ç†PDFæ–‡ä»¶ - ä¼˜å…ˆæå–æ–‡æœ¬ï¼Œå¤±è´¥åˆ™è¿”å›äºŒè¿›åˆ¶"""
        try:
            # é¦–å…ˆå°è¯•æå–æ–‡æœ¬
            try:
                import PyPDF2
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text_parts = []
                    
                    # è¯»å–å‰10é¡µ
                    max_pages = min(10, len(reader.pages))
                    for page_num in range(max_pages):
                        page = reader.pages[page_num]
                        text_parts.append(page.extract_text())
                    
                    result['text_content'] = '\n'.join(text_parts)
                    result['metadata']['pages'] = len(reader.pages)
                    result['metadata']['extracted_pages'] = max_pages
                    
                    if result['text_content'].strip():
                        result['success'] = True
                        print(f"[FileProcessor] æˆåŠŸä»PDFæå–æ–‡æœ¬: {result['filename']} ({max_pages}é¡µ)")
                        return result
                    else:
                        # æå–çš„æ–‡æœ¬ä¸ºç©ºï¼Œä½¿ç”¨äºŒè¿›åˆ¶
                        raise ValueError("PDFæ–‡æœ¬ä¸ºç©º")
                        
            except (ImportError, Exception) as e:
                # PyPDF2ä¸å¯ç”¨æˆ–æå–å¤±è´¥ï¼Œè¯»å–äºŒè¿›åˆ¶
                print(f"[FileProcessor] PDFæ–‡æœ¬æå–å¤±è´¥ï¼Œä½¿ç”¨äºŒè¿›åˆ¶æ¨¡å¼: {e}")
                with open(filepath, 'rb') as f:
                    result['binary_data'] = f.read()
                result['success'] = True
                return result
                
        except Exception as e:
            result['error'] = f"è¯»å–PDFå¤±è´¥: {str(e)}"
            return result
    
    @staticmethod
    def _process_word(filepath: str, result: Dict) -> Dict:
        """å¤„ç†Wordæ–‡æ¡£"""
        try:
            try:
                from docx import Document
                doc = Document(filepath)
                
                # æå–æ‰€æœ‰æ®µè½æ–‡æœ¬
                paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
                result['text_content'] = '\n'.join(paragraphs)
                
                # æå–è¡¨æ ¼å†…å®¹
                if doc.tables:
                    result['text_content'] += '\n\n[è¡¨æ ¼å†…å®¹]:\n'
                    for table in doc.tables:
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            result['text_content'] += ' | '.join(cells) + '\n'
                
                result['metadata']['paragraphs'] = len(paragraphs)
                result['metadata']['tables'] = len(doc.tables)
                result['metadata']['chars'] = len(result['text_content'])
                
                result['success'] = True
                print(f"[FileProcessor] æˆåŠŸæå–Wordæ–‡æ¡£: {result['filename']} ({len(result['text_content'])} å­—ç¬¦)")
                return result
                
            except ImportError:
                result['error'] = "æœªå®‰è£…python-docxåº“ï¼Œæ— æ³•è¯»å–Wordæ–‡æ¡£ã€‚è¯·è¿è¡Œ: pip install python-docx"
                return result
                
        except Exception as e:
            result['error'] = f"è¯»å–Wordæ–‡æ¡£å¤±è´¥: {str(e)}"
            return result

    @staticmethod
    def _process_powerpoint(filepath: str, result: Dict) -> Dict:
        """å¤„ç†PowerPointæ–‡æ¡£ï¼ˆä¼˜å…ˆ .pptxï¼‰"""
        try:
            from pptx import Presentation

            presentation = Presentation(filepath)
            slide_texts = []

            for idx, slide in enumerate(presentation.slides, start=1):
                parts = []
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text:
                            text = shape.text.strip()
                            if text:
                                parts.append(text)
                    except Exception:
                        continue

                if parts:
                    slide_texts.append(f"[Slide {idx}]\n" + "\n".join(parts))

            result['text_content'] = "\n\n".join(slide_texts)
            result['metadata']['slides'] = len(presentation.slides)
            result['metadata']['extracted_slides'] = len(slide_texts)

            if result['text_content'].strip():
                result['success'] = True
                print(f"[FileProcessor] æˆåŠŸæå–PowerPointæ–‡æ¡£: {result['filename']} ({len(result['text_content'])} å­—ç¬¦)")
            else:
                # æ— æ–‡æœ¬æ—¶ä»æ ‡è®°æˆåŠŸï¼Œé¿å…åç»­æŠŠå®ƒå½“æˆé”™è¯¯æ–‡ä»¶
                result['success'] = True
                print(f"[FileProcessor] PowerPointæå–å®Œæˆä½†æ— å¯è¯»æ–‡æœ¬: {result['filename']}")

            return result

        except Exception as e:
            result['error'] = f"è¯»å–PowerPointæ–‡æ¡£å¤±è´¥: {str(e)}"
            return result
    
    @staticmethod
    def _process_excel(filepath: str, result: Dict) -> Dict:
        """å¤„ç†Excelæ–‡ä»¶"""
        try:
            try:
                import pandas as pd
                
                # è¯»å–Excelæ–‡ä»¶
                xls = pd.ExcelFile(filepath)
                text_parts = [f"Excelå·¥ä½œç°¿: {result['filename']}"]
                text_parts.append(f"å·¥ä½œè¡¨: {', '.join(xls.sheet_names)}\n")
                
                # è¯»å–å‰3ä¸ªå·¥ä½œè¡¨
                for sheet_name in xls.sheet_names[:3]:
                    df = pd.read_excel(filepath, sheet_name=sheet_name)
                    text_parts.append(f"\n=== å·¥ä½œè¡¨: {sheet_name} ===")
                    text_parts.append(df.to_string(max_rows=50, max_cols=10))
                
                result['text_content'] = '\n'.join(text_parts)
                result['metadata']['sheets'] = len(xls.sheet_names)
                result['metadata']['extracted_sheets'] = min(3, len(xls.sheet_names))
                
                result['success'] = True
                print(f"[FileProcessor] æˆåŠŸæå–Excelæ•°æ®: {result['filename']} ({len(xls.sheet_names)} ä¸ªå·¥ä½œè¡¨)")
                return result
                
            except ImportError:
                result['error'] = "æœªå®‰è£…pandas/openpyxlåº“ï¼Œæ— æ³•è¯»å–Excelæ–‡ä»¶ã€‚è¯·è¿è¡Œ: pip install pandas openpyxl"
                return result
                
        except Exception as e:
            result['error'] = f"è¯»å–Excelæ–‡ä»¶å¤±è´¥: {str(e)}"
            return result
    
    @staticmethod
    def _process_text(filepath: str, result: Dict) -> Dict:
        """å¤„ç†æ–‡æœ¬æ–‡ä»¶ - å°è¯•å¤šç§ç¼–ç """
        encodings = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as f:
                    result['text_content'] = f.read()
                
                result['metadata']['encoding'] = encoding
                result['metadata']['lines'] = result['text_content'].count('\n') + 1
                result['metadata']['chars'] = len(result['text_content'])
                
                result['success'] = True
                print(f"[FileProcessor] æˆåŠŸè¯»å–æ–‡æœ¬æ–‡ä»¶: {result['filename']} (ç¼–ç : {encoding})")
                return result
                
            except (UnicodeDecodeError, LookupError):
                continue
        
        # æ‰€æœ‰ç¼–ç éƒ½å¤±è´¥
        result['error'] = f"æ— æ³•è§£ç æ–‡ä»¶ï¼Œå°è¯•è¿‡çš„ç¼–ç : {', '.join(encodings)}"
        return result
    
    @staticmethod
    def format_result_for_chat(result: Dict, user_message: str = "") -> Tuple[str, Optional[Dict]]:
        """
        æ ¼å¼åŒ–å¤„ç†ç»“æœï¼Œç”¨äºå‘é€ç»™AIèŠå¤©
        
        Returns:
            (formatted_message, file_data)
            formatted_message: åŒ…å«æ–‡ä»¶å†…å®¹çš„æ¶ˆæ¯
            file_data: å¦‚æœæ˜¯äºŒè¿›åˆ¶æ–‡ä»¶ï¼ˆå›¾ç‰‡/PDFï¼‰ï¼Œè¿”å› {'mime_type': str, 'data': bytes}
        """
        if not result['success']:
            error_msg = f"âŒ æ–‡ä»¶å¤„ç†å¤±è´¥: {result['error']}"
            return f"{user_message}\n\n{error_msg}", None
        
        # å¦‚æœæœ‰äºŒè¿›åˆ¶æ•°æ®ï¼ˆå›¾ç‰‡/PDFï¼‰
        if result['binary_data']:
            message = user_message
            file_data = {
                'mime_type': result['mime_type'],
                'data': result['binary_data']
            }
            return message, file_data
        
        # å¦‚æœæœ‰æ–‡æœ¬å†…å®¹
        if result['text_content']:
            metadata_str = ""
            if result['metadata']:
                meta_items = [f"{k}: {v}" for k, v in result['metadata'].items()]
                metadata_str = f" ({', '.join(meta_items)})"
            
            formatted = f"{user_message}\n\n"
            formatted += f"ğŸ“„ æ–‡ä»¶: {result['filename']}{metadata_str}\n\n"
            formatted += f"=== æ–‡ä»¶å†…å®¹ ===\n{result['text_content']}"
            
            return formatted, None
        
        # æœªçŸ¥æƒ…å†µ
        return f"{user_message}\n\nâš ï¸ æ–‡ä»¶å·²ä¸Šä¼ ä½†æœªæå–åˆ°å†…å®¹", None


# ä¾¿æ·å‡½æ•°
def process_uploaded_file(filepath: str, user_message: str = "") -> Tuple[str, Optional[Dict]]:
    """
    å¤„ç†ä¸Šä¼ çš„æ–‡ä»¶ï¼Œè¿”å›æ ¼å¼åŒ–çš„æ¶ˆæ¯å’Œæ–‡ä»¶æ•°æ®
    
    Args:
        filepath: æ–‡ä»¶è·¯å¾„
        user_message: ç”¨æˆ·çš„æ¶ˆæ¯
        
    Returns:
        (formatted_message, file_data)
    """
    processor = FileProcessor()
    result = processor.process_file(filepath)
    return processor.format_result_for_chat(result, user_message)
