#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文档智能编辑器 - AI驱动的文档修改系统
读取文档 → AI分析 → 应用修改
"""

import os
import json
from typing import Dict, List, Any, Optional
from datetime import datetime


class DocumentEditor:
    """文档智能编辑器"""
    
    def __init__(self):
        pass
    
    @staticmethod
    def edit_ppt(file_path: str, modifications: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        根据AI建议修改PPT
        
        Args:
            file_path: PPT文件路径
            modifications: 修改指令列表
                [
                    {
                        "slide_index": 0,
                        "action": "update_title" | "update_content" | "add_content" | "delete_content",
                        "target": "title" | "content" | "notes",
                        "content": "新内容",
                        "position": 0  # 对于列表项
                    },
                    ...
                ]
        
        Returns:
            修改结果
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt
            
            # 读取原文件
            prs = Presentation(file_path)
            
            applied_count = 0
            errors = []
            
            for mod in modifications:
                try:
                    slide_index = mod.get("slide_index")
                    action = mod.get("action")
                    target = mod.get("target", "content")
                    content = mod.get("content", "")
                    
                    if slide_index >= len(prs.slides):
                        errors.append(f"幻灯片索引{slide_index}超出范围")
                        continue
                    
                    slide = prs.slides[slide_index]
                    
                    # 修改标题
                    if action == "update_title" and target == "title":
                        if slide.shapes.title:
                            slide.shapes.title.text = content
                            applied_count += 1
                    
                    # 修改内容
                    elif action == "update_content" and target == "content":
                        position = mod.get("position", 0)
                        # 找到第一个非标题的文本框
                        content_shape = None
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                                content_shape = shape
                                break
                        
                        if content_shape:
                            paragraphs = content_shape.text_frame.paragraphs
                            if position < len(paragraphs):
                                paragraphs[position].text = content
                                applied_count += 1
                    
                    # 添加内容
                    elif action == "add_content" and target == "content":
                        # 找到或创建内容文本框
                        content_shape = None
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                                content_shape = shape
                                break
                        
                        if content_shape:
                            tf = content_shape.text_frame
                            p = tf.add_paragraph()
                            p.text = content
                            p.level = 0
                            applied_count += 1
                    
                    # 删除内容
                    elif action == "delete_content" and target == "content":
                        position = mod.get("position", 0)
                        content_shape = None
                        for shape in slide.shapes:
                            if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                                content_shape = shape
                                break
                        
                        if content_shape:
                            paragraphs = content_shape.text_frame.paragraphs
                            if position < len(paragraphs):
                                # 清空段落内容（python-pptx不支持直接删除段落）
                                paragraphs[position].text = ""
                                applied_count += 1
                    
                    # 修改备注
                    elif target == "notes":
                        if slide.has_notes_slide:
                            notes_slide = slide.notes_slide
                            notes_slide.notes_text_frame.text = content
                        else:
                            # 创建备注
                            notes_slide = slide.notes_slide
                            notes_slide.notes_text_frame.text = content
                        applied_count += 1
                
                except Exception as e:
                    errors.append(f"应用修改失败: {str(e)}")
            
            # 保存到新文件
            base_name = os.path.splitext(file_path)[0]
            new_file_path = f"{base_name}_edited_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(new_file_path)
            
            return {
                "success": True,
                "file_path": new_file_path,
                "applied_count": applied_count,
                "total_modifications": len(modifications),
                "errors": errors
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"编辑PPT失败: {str(e)}"
            }
    
    @staticmethod
    def edit_word(file_path: str, modifications: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        根据AI建议修改Word文档
        
        Args:
            modifications: 修改指令
                [
                    {
                        "paragraph_index": 0,
                        "action": "update" | "insert" | "delete",
                        "content": "新内容"
                    },
                    ...
                ]
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            applied_count = 0
            errors = []
            
            for mod in modifications:
                try:
                    para_index = mod.get("paragraph_index")
                    action = mod.get("action")
                    content = mod.get("content", "")
                    
                    if action == "update":
                        if para_index < len(doc.paragraphs):
                            doc.paragraphs[para_index].text = content
                            applied_count += 1
                    
                    elif action == "insert":
                        # 在指定位置插入新段落
                        p = doc.add_paragraph(content)
                        if para_index < len(doc.paragraphs):
                            # 移动到指定位置
                            p._element.getparent().remove(p._element)
                            doc.paragraphs[para_index]._element.addprevious(p._element)
                        applied_count += 1
                    
                    elif action == "delete":
                        if para_index < len(doc.paragraphs):
                            p = doc.paragraphs[para_index]
                            p._element.getparent().remove(p._element)
                            applied_count += 1
                
                except Exception as e:
                    errors.append(f"应用修改失败: {str(e)}")
            
            # 保存
            base_name = os.path.splitext(file_path)[0]
            new_file_path = f"{base_name}_edited_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            doc.save(new_file_path)
            
            return {
                "success": True,
                "file_path": new_file_path,
                "applied_count": applied_count,
                "total_modifications": len(modifications),
                "errors": errors
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"编辑Word失败: {str(e)}"
            }
    
    @staticmethod
    def edit_excel(file_path: str, modifications: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        根据AI建议修改Excel
        
        Args:
            modifications: 修改指令
                [
                    {
                        "sheet_name": "Sheet1",
                        "row": 0,
                        "col": 0,
                        "action": "update" | "insert_row" | "delete_row",
                        "value": "新值"
                    },
                    ...
                ]
        """
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path)
            applied_count = 0
            errors = []
            
            for mod in modifications:
                try:
                    sheet_name = mod.get("sheet_name")
                    action = mod.get("action")
                    
                    if sheet_name not in wb.sheetnames:
                        errors.append(f"工作表'{sheet_name}'不存在")
                        continue
                    
                    sheet = wb[sheet_name]
                    
                    if action == "update":
                        row = mod.get("row", 0) + 1  # Excel行从1开始
                        col = mod.get("col", 0) + 1
                        value = mod.get("value", "")
                        sheet.cell(row=row, column=col, value=value)
                        applied_count += 1
                    
                    elif action == "insert_row":
                        row = mod.get("row", 0) + 1
                        sheet.insert_rows(row)
                        applied_count += 1
                    
                    elif action == "delete_row":
                        row = mod.get("row", 0) + 1
                        sheet.delete_rows(row)
                        applied_count += 1
                
                except Exception as e:
                    errors.append(f"应用修改失败: {str(e)}")
            
            # 保存
            base_name = os.path.splitext(file_path)[0]
            new_file_path = f"{base_name}_edited_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
            wb.save(new_file_path)
            
            return {
                "success": True,
                "file_path": new_file_path,
                "applied_count": applied_count,
                "total_modifications": len(modifications),
                "errors": errors
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"编辑Excel失败: {str(e)}"
            }
    
    @staticmethod
    def parse_ai_suggestions(ai_response: str) -> List[Dict[str, Any]]:
        """
        解析AI返回的修改建议
        
        期望AI返回JSON格式:
        ```json
        {
            "modifications": [
                {
                    "slide_index": 0,
                    "action": "update_title",
                    "target": "title",
                    "content": "新标题"
                },
                ...
            ]
        }
        ```
        
        Returns:
            修改指令列表
        """
        try:
            # 尝试提取JSON（AI可能返回包含其他文本的响应）
            import re
            
            # 查找JSON块
            json_match = re.search(r'```json\s*(\{.*?\})\s*```', ai_response, re.DOTALL)
            if json_match:
                json_str = json_match.group(1)
            else:
                # 尝试查找裸JSON
                json_match = re.search(r'\{.*"modifications".*\}', ai_response, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                else:
                    json_str = ai_response
            
            data = json.loads(json_str)
            
            if isinstance(data, dict) and "modifications" in data:
                return data["modifications"]
            elif isinstance(data, list):
                return data
            else:
                return []
        
        except Exception as e:
            print(f"[DocumentEditor] 解析AI建议失败: {e}")
            return []


    @staticmethod
    def edit_word_with_tracked_changes(file_path: str, annotations: List[Dict[str, str]]) -> Dict[str, Any]:
        """
        直接编辑Word文档，并将修改后的文本用绿色标记
        使用更可靠的XML级别操作来处理复杂的段落结构
        
        Args:
            file_path: Word文档路径  
            annotations: 标注列表，每个包含：
                - 原文片段
                - 修改后文本
                - 修改建议
                - 理由
        
        Returns:
            编辑结果统计
        """
        try:
            from docx import Document
            from docx.shared import RGBColor
            
            doc = Document(file_path)
            
            applied_count = 0
            failed_count = 0
            
            print(f"[Editor] 📝 开始应用编辑修改...")
            print(f"[Editor] 📊 收到 {len(annotations)} 条标注")
            
            # Debug: 检查第一条标注的结构
            if annotations:
                first_anno = annotations[0]
                print(f"[Editor] 🔍 第一条标注示例: {first_anno}")
            
            # 对每个标注进行处理
            for anno in annotations:
                original = anno.get("原文片段", "").strip()
                modified = anno.get("修改后文本", "").strip()
                
                if not original or not modified:
                    continue
                
                # 在所有段落中查找并替换
                found = False
                for para in doc.paragraphs:
                    # 检查这个段落中是否包含原文
                    if original not in para.text:
                        continue
                    
                    # 找到了，执行替换
                    # 方法：完全重建段落文本
                    para_text = para.text
                    new_text = para_text.replace(original, original + "<<<MODIFIED>>>" + modified + "<<<END>>>")
                    
                    # 清空并重新构建段落
                    # 获取原样式
                    style = para.style
                    paragraph_format = para.paragraph_format
                    
                    # 提取所有原始runs的属性信息
                    original_runs_props = []
                    for run in para.runs:
                        props = {
                            'text': run.text,
                            'bold': run.font.bold,
                            'italic': run.font.italic,
                            'size': run.font.size,
                        }
                        original_runs_props.append(props)
                    
                    # 清空paragraph中的所有runs
                    for run in list(para.runs):
                        r_element = run._element
                        r_element.getparent().remove(r_element)
                    
                    # 重新添加文本，处理修改部分
                    parts = new_text.split("<<<MODIFIED>>>")
                    for i, part in enumerate(parts):
                        if i == 0:
                            # 第一部分是未修改的文本
                            run = para.add_run(part)
                        else:
                            # 分离修改前和修改后
                            sub_parts = part.split("<<<END>>>")
                            if len(sub_parts) >= 2:
                                # 添加修改后的部分（绿色+加粗）
                                run = para.add_run(sub_parts[0])
                                run.font.color.rgb = RGBColor(0, 128, 0)  # 绿色
                                run.font.bold = True
                                
                                # 添加后续的文本
                                if sub_parts[1]:
                                    para.add_run(sub_parts[1])
                    
                    applied_count += 1
                    found = True
                    print(f"  ✅ 修改: '{original}' → '{modified}'")
                    break
                
                if not found:
                    failed_count += 1
                    print(f"  ⚠️ 未找到: '{original}'")
            
            # 生成修改版文件名
            if file_path.endswith('.docx'):
                revised_file = file_path.replace('.docx', '_revised.docx')
            else:
                revised_file = file_path + '_revised.docx'
            
            # 保存文档
            doc.save(revised_file)
            print(f"[Editor] 💾 修改版已保存 ({applied_count}个修改)")
            
            return {
                "success": True,
                "applied": applied_count,
                "failed": failed_count,
                "total": len(annotations),
                "file_path": revised_file
            }
            
        except Exception as e:
            print(f"[Editor] ❌ 编辑失败: {str(e)}")
            import traceback
            traceback.print_exc()
            return {
                "success": False,
                "error": str(e)
            }


if __name__ == "__main__":
    # 测试示例
    editor = DocumentEditor()
    
    # 示例：修改PPT
    sample_modifications = [
        {
            "slide_index": 0,
            "action": "update_title",
            "target": "title",
            "content": "修改后的标题"
        },
        {
            "slide_index": 1,
            "action": "add_content",
            "target": "content",
            "content": "AI建议添加的新要点"
        }
    ]
    
    print("文档编辑器准备就绪")
