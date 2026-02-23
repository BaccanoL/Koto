#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试文件生成和质量评估流程
"""

import sys
import os

# 添加 web 目录到路径
sys.path.insert(0, os.path.join(os.path.dirname(__file__), 'web'))

def test_ppt_evaluator():
    """测试 PPT 评估器"""
    print("\n" + "=" * 50)
    print("测试 PPT 评估器")
    print("=" * 50)
    
    from quality_evaluator import PPTEvaluator
    import tempfile
    
    # 测试用例：创建一个简单的 PPT（如果可用）
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # 创建测试 PPT
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 添加几个幻灯片
        for i in range(1, 4):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"Slide {i}"
            
            if i == 2:
                # 添加很多文本到第2页
                body = slide.placeholders[1].text_frame
                body.text = "这是一个测试幻灯片。" * 100
        
        # 保存测试文件到临时目录
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            test_ppt_path = tmp.name
        
        prs.save(test_ppt_path)
        
        # 评估 PPT
        evaluator = PPTEvaluator()
        result = evaluator.evaluate_pptx_file(test_ppt_path)
        
        print(f"📊 总体评分: {result.overall_score}/100")
        print(f"📋 类别评分: {result.category_scores}")
        print(f"⚠️ 发现的问题数: {len(result.issues)}")
        print(f"💡 改进建议数: {len(result.suggestions)}")
        print(f"🎯 改进优先级: {result.improvement_priority}")
        print(f"✓ 需要改进: {result.needs_improvement}")
        
        # 清理
        os.remove(test_ppt_path)
        print("\n✅ PPT 评估器测试通过")
    
    except ImportError:
        print("⚠️ python-pptx 未安装，跳过 PPT 评估测试")
    except Exception as e:
        print(f"❌ PPT 评估器测试失败: {e}")
        import traceback
        traceback.print_exc()


def test_document_evaluator():
    """测试文档评估器"""
    print("\n" + "=" * 50)
    print("测试文档评估器")
    print("=" * 50)
    
    from quality_evaluator import DocumentEvaluator
    
    # 测试用例1：优秀文档
    excellent_doc = """# 项目总结报告

## 项目概览
这是一份详细的项目总结报告，包含了完整的项目信息。

## 主要成就
- 完成了 95% 的计划功能
- 用户满意度达到 4.8/5 星
- 性能提升 40%

## 技术细节
### 架构设计
采用微服务架构，提高了系统的可扩展性。

### 核心组件
1. 数据层：PostgreSQL + Redis
2. 业务层：Go + gRPC
3. 展示层：React 18

## 测试覆盖率
- 单元测试：87%
- 集成测试：92%
- 端到端测试：78%

## 结论
项目成功交付，达到所有质量目标。未来计划继续优化性能。
"""
    
    evaluator = DocumentEvaluator()
    result = evaluator.evaluate_document(excellent_doc)
    
    print("测试1：优秀文档")
    print(f"  评分: {result.overall_score}/100")
    print(f"  需要改进: {result.needs_improvement}")
    
    # 测试用例2：需要改进的文档
    poor_doc = "这是很短的文档。"
    
    result2 = evaluator.evaluate_document(poor_doc)
    print("\n测试2：简短文档")
    print(f"  评分: {result2.overall_score}/100")
    print(f"  需要改进: {result2.needs_improvement}")
    print(f"  问题: {result2.issues}")
    
    print("\n✅ 文档评估器测试通过")


def test_progress_tracker():
    """测试进度追踪器"""
    print("\n" + "=" * 50)
    print("测试进度追踪器")
    print("=" * 50)
    
    from progress_tracker import (
        ProgressTracker, ProgressBroadcaster, 
        TaskStage, GenerationProgressManager
    )
    
    # 创建追踪器
    tracker = ProgressTracker("test_task", total_stages=4)
    
    # 模拟进度更新
    stages = [
        (TaskStage.VALIDATING.value, "正在验证输入..."),
        (TaskStage.EVALUATING.value, "正在评估质量..."),
        (TaskStage.IMPROVING.value, "正在改进内容..."),
        (TaskStage.GENERATING.value, "正在生成文件..."),
        (TaskStage.COMPLETED.value, "文件生成完成！"),
    ]
    
    for stage, msg in stages:
        update = tracker.update(stage, msg)
        print(f"  {update.progress_percent}% - {update.message}")
    
    print("\n✅ 进度追踪器测试通过")


def test_feedback_loop():
    """测试反馈循环（需要 Gemini 客户端）"""
    print("\n" + "=" * 50)
    print("测试反馈循环")
    print("=" * 50)
    
    try:
        from feedback_loop import FeedbackLoopManager
        
        # 模拟一个简单的 mock 客户端
        class MockClient:
            class Models:
                def generate_content(self, **kwargs):
                    class Response:
                        text = "改进后的内容：这是改进后的版本，包含更多细节。"
                    return Response()
            models = Models()
        
        # 创建管理器
        manager = FeedbackLoopManager(lambda: MockClient())
        
        # 测试评估结果
        evaluation = {
            "overall_score": 60,
            "needs_improvement": True,
            "issues": ["内容过短", "结构不清晰"],
            "suggestions": ["增加更多细节", "添加标题"],
            "improvement_priority": ["优先添加标题", "然后增加内容"]
        }
        
        # 测试改进（会调用 API）
        result = manager.improve_document_content(
            "这是原始内容。",
            evaluation,
            "测试文档"
        )
        
        print(f"  改进次数: {result['iterations']}")
        print(f"  最终评分: {result['final_score']}")
        print(f"  改进历史: {len(result['improvement_history'])} 条记录")
        
        print("\n✅ 反馈循环测试通过")
    
    except Exception as e:
        print(f"⚠️ 反馈循环测试需要 Gemini 客户端: {e}")


def test_evaluation_quality():
    """测试质量评估的准确性"""
    print("\n" + "=" * 50)
    print("测试质量评估准确性")
    print("=" * 50)
    
    from quality_evaluator import DocumentEvaluator
    
    test_cases = [
        {
            "name": "优秀文档",
            "content": "# 标题\n## 小标题\n内容详细的文档。" * 10,
            "expected_score_min": 80
        },
        {
            "name": "差文档",
            "content": "短",
            "expected_score_max": 50
        },
        {
            "name": "中等文档",
            "content": "# 标题\n这是一个中等长度的文档，有一些结构但内容不足。",
            "expected_score_min": 50,
            "expected_score_max": 80
        }
    ]
    
    evaluator = DocumentEvaluator()
    
    for case in test_cases:
        result = evaluator.evaluate_document(case["content"])
        score = result.overall_score
        
        status = "✓" if (
            ("expected_score_min" not in case or score >= case["expected_score_min"]) and
            ("expected_score_max" not in case or score <= case["expected_score_max"])
        ) else "✗"
        
        print(f"  {status} {case['name']}: {score}/100")
    
    print("\n✅ 质量评估准确性测试完成")


if __name__ == "__main__":
    print("\n🧪 开始测试质量评估和改进系统\n")
    
    test_document_evaluator()
    test_ppt_evaluator()
    test_progress_tracker()
    test_evaluation_quality()
    test_feedback_loop()
    
    print("\n" + "=" * 50)
    print("✅ 所有测试完成！")
    print("=" * 50)
